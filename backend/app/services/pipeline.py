# [ROLE] 文書インデックス化パイプライン（抽出→Markdown化→チャンキング→品質フィルタ→埋め込み→ChromaDB格納）
# [DEPS] core/config.py, core/retry.py, services/embedder.py, models/db.py
# [CALLED_BY] routers/documents.py

import asyncio
import re
from typing import Optional, TypedDict
from uuid import uuid4

import chromadb
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from langchain.text_splitter import (
    MarkdownHeaderTextSplitter,
    RecursiveCharacterTextSplitter,
)
from pptx import Presentation

from app.core.config import settings
from app.core.retry import with_retry
from app.models.db import Document, OcrResult, SessionLocal
from app.services.embedder import get_embedder


HEADERS_TO_SPLIT = [
    ("#", "chapter"),
    ("##", "section"),
    ("###", "subsection"),
]

_PAGE_MARKER_RE = re.compile(r"<!--\s*page:\s*(\d+)\s*-->")
_SLIDE_MARKER_RE = re.compile(r"<!--\s*slide:\s*(\d+)\s*-->")
_CHAPTER_RE = re.compile(r"^第[\d０-９一二三四五六七八九十百]+章[ \t　]+.+$")
_SECTION_RE = re.compile(r"^(\d+\.\d+)[ \t　]+(.+)$")
_SUBSECTION_RE = re.compile(r"^(\d+\.\d+\.\d+)[ \t　]+(.+)$")


_ocr = None


def get_ocr():
    """PaddleOCR モデルのシングルトン。初回呼び出し時にロードしてグローバルにキャッシュする。"""
    global _ocr
    if _ocr is None:
        from paddleocr import PaddleOCR
        _ocr = PaddleOCR(use_angle_cls=True, lang="japan")
    return _ocr


class ChunkMetadata(TypedDict):
    doc_id: str
    source_type: str
    title: str
    author: str
    year: int
    page_num: int
    chapter: str
    section: str


# ==================== Step 1: テキスト抽出 ====================

def extract_pdf(file_path: str) -> list[dict]:
    """PDF からページ単位でテキスト抽出。戻り値: [{"page": 1, "text": "..."}]"""
    doc = fitz.open(file_path)
    try:
        return [
            {"page": i + 1, "text": page.get_text()}
            for i, page in enumerate(doc)
            if page.get_text().strip()
        ]
    finally:
        doc.close()


def extract_pptx(file_path: str) -> list[dict]:
    """PPTX からスライド単位でテキスト抽出。スライド本文＋話者ノートを連結。"""
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        texts.append(line)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                texts.append(notes)
        body = "\n".join(texts)
        if body.strip():
            slides.append({"slide": i + 1, "text": body})
    return slides


def extract_docx(file_path: str) -> list[dict]:
    """DOCX から段落単位で抽出。ページ概念がないため概算で 1 ページ扱い。"""
    doc = DocxDocument(file_path)
    body = "\n".join(p.text for p in doc.paragraphs if p.text and p.text.strip())
    if not body.strip():
        return []
    return [{"page": 1, "text": body}]


def extract_ocr(file_path: str) -> dict:
    """
    画像 OCR。テキストブロックごとに信頼度スコアとバウンディングボックスを取得する。
    blocks には全ブロック（低信頼度含む）を保持し、
    インデックス化時の除外は to_markdown_ocr 側で行う。
    avg_confidence は全ブロックの平均、low_conf_count は閾値未満ブロック数。
    """
    ocr = get_ocr()
    result = ocr.ocr(file_path, cls=True)
    page_result = result[0] if result and result[0] else []

    blocks: list[dict] = []
    scores: list[float] = []
    low_conf_count = 0
    threshold = settings.ocr_confidence_threshold

    for det in page_result:
        try:
            bbox_pts, (text, score) = det
        except (TypeError, ValueError):
            continue
        score_f = float(score)
        scores.append(score_f)
        if score_f < threshold:
            low_conf_count += 1
        # 4頂点 [[x1,y1],[x2,y1],[x2,y2],[x1,y2]] を [xmin, ymin, xmax, ymax] に変換
        try:
            xs = [float(pt[0]) for pt in bbox_pts]
            ys = [float(pt[1]) for pt in bbox_pts]
            bbox = [min(xs), min(ys), max(xs), max(ys)]
        except (TypeError, ValueError, IndexError):
            bbox = [0.0, 0.0, 0.0, 0.0]
        blocks.append({"text": text, "confidence": score_f, "bbox": bbox})

    avg_confidence = (sum(scores) / len(scores)) if scores else 0.0
    return {
        "blocks": blocks,
        "avg_confidence": avg_confidence,
        "low_conf_count": low_conf_count,
    }


# ==================== Step 2: Markdown 構造化 ====================

def _detect_heading(line: str) -> Optional[str]:
    """1 行を見出しに変換。マッチしなければ None を返す。"""
    s = line.rstrip()
    if _SUBSECTION_RE.match(s):
        m = _SUBSECTION_RE.match(s)
        return f"### {m.group(1)} {m.group(2)}"
    if _SECTION_RE.match(s):
        m = _SECTION_RE.match(s)
        return f"## {m.group(1)} {m.group(2)}"
    if _CHAPTER_RE.match(s):
        return f"# {s}"
    return None


def _structure_text(text: str) -> str:
    """章・節パターンを Markdown 見出しに変換する。"""
    out: list[str] = []
    for raw in text.splitlines():
        heading = _detect_heading(raw)
        if heading:
            out.append("")
            out.append(heading)
            out.append("")
        else:
            out.append(raw)
    return "\n".join(out)


def to_markdown_pdf(pages: list[dict]) -> str:
    """PDF のページリストを Markdown 化。各ページ末に <!-- page: N --> を挿入。"""
    parts: list[str] = []
    for p in pages:
        parts.append(_structure_text(p["text"]))
        parts.append(f"<!-- page: {p['page']} -->")
    return "\n".join(parts)


def to_markdown_pptx(slides: list[dict]) -> str:
    """PPTX のスライドリストを Markdown 化。各スライド末に <!-- slide: N --> を挿入。"""
    parts: list[str] = []
    for s in slides:
        parts.append(_structure_text(s["text"]))
        parts.append(f"<!-- slide: {s['slide']} -->")
    return "\n".join(parts)


def to_markdown_docx(pages: list[dict]) -> str:
    """DOCX を Markdown 化。"""
    parts: list[str] = []
    for p in pages:
        parts.append(_structure_text(p["text"]))
        parts.append(f"<!-- page: {p['page']} -->")
    return "\n".join(parts)


def to_markdown_ocr(ocr_result: dict) -> str:
    """
    OCR 抽出結果を Markdown 化。インデックス化対象は信頼度 >= 閾値のブロックのみ。
    （低信頼度ブロックは ocr_results テーブルには残すが索引対象からは外す）
    """
    threshold = settings.ocr_confidence_threshold
    text = "\n".join(
        b["text"] for b in ocr_result["blocks"] if b["confidence"] >= threshold
    )
    if not text:
        return "<!-- page: 1 -->"
    return _structure_text(text) + "\n<!-- page: 1 -->"


# ==================== Step 4: ハイブリッドチャンキング ====================

def chunk_markdown(markdown_text: str) -> list:
    """Markdown ヘッダ分割 → 文字数分割のハイブリッドチャンキング。"""
    md_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=HEADERS_TO_SPLIT)
    md_chunks = md_splitter.split_text(markdown_text)
    char_splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
    )
    return char_splitter.split_documents(md_chunks)


# ==================== Step 6: 埋め込み生成 ====================

def embed_chunks(chunks: list) -> list[list[float]]:
    """チャンク本文に 'passage: ' プレフィックスを付けて埋め込みを生成。"""
    texts = ["passage: " + c.page_content for c in chunks]
    embeddings = get_embedder().encode(texts)
    if hasattr(embeddings, "tolist"):
        return embeddings.tolist()
    return list(embeddings)


# ==================== Step 7: ChromaDB 格納 ====================

def store_to_chroma(chunks: list, embeddings: list[list[float]], metadatas: list[dict]) -> None:
    client = chromadb.PersistentClient(path=settings.chroma_path)
    collection = client.get_or_create_collection(settings.chroma_collection)
    collection.add(
        embeddings=embeddings,
        documents=[c.page_content for c in chunks],
        metadatas=metadatas,
        ids=[str(uuid4()) for _ in chunks],
    )


# ==================== ヘルパ ====================

def _extract_page_num_from_chunk(content: str) -> int:
    m = _PAGE_MARKER_RE.search(content)
    if m:
        return int(m.group(1))
    m = _SLIDE_MARKER_RE.search(content)
    if m:
        return int(m.group(1))
    return 0


def _strip_page_markers(content: str) -> str:
    content = _PAGE_MARKER_RE.sub("", content)
    content = _SLIDE_MARKER_RE.sub("", content)
    return content.strip()


def _build_metadata(
    chunk,
    doc_id: str,
    source_type: str,
    title: str,
    author: str,
    year: int,
) -> ChunkMetadata:
    md = chunk.metadata or {}
    return {
        "doc_id": doc_id,
        "source_type": source_type,
        "title": title or "",
        "author": author or "",
        "year": int(year) if year else 0,
        "page_num": _extract_page_num_from_chunk(chunk.page_content),
        "chapter": md.get("chapter", "") or "",
        "section": md.get("section", "") or "",
    }


def _extract_to_markdown(source_type: str, file_path: str) -> tuple[str, Optional[dict]]:
    """
    戻り値: (markdown_text, ocr_result)
    ocr_result は source_type が 'ocr' の場合のみ dict、それ以外は None。
    """
    if source_type == "pdf":
        return to_markdown_pdf(extract_pdf(file_path)), None
    if source_type == "pptx":
        return to_markdown_pptx(extract_pptx(file_path)), None
    if source_type == "word":
        return to_markdown_docx(extract_docx(file_path)), None
    if source_type == "ocr":
        ocr_result = extract_ocr(file_path)
        return to_markdown_ocr(ocr_result), ocr_result
    raise ValueError(f"Unsupported source_type: {source_type}")


def _save_ocr_result(doc_id: str, ocr_result: dict) -> None:
    """OCR の信頼度・バウンディングボックスを ocr_results テーブルへ upsert する。
    再実行時の unique 衝突を避けるため、既存レコードがあれば更新する。"""
    db = SessionLocal()
    try:
        existing = (
            db.query(OcrResult).filter(OcrResult.document_id == doc_id).first()
        )
        if existing:
            existing.avg_confidence = ocr_result["avg_confidence"]
            existing.low_conf_count = ocr_result["low_conf_count"]
            existing.blocks = ocr_result["blocks"]
        else:
            db.add(
                OcrResult(
                    document_id=doc_id,
                    avg_confidence=ocr_result["avg_confidence"],
                    low_conf_count=ocr_result["low_conf_count"],
                    blocks=ocr_result["blocks"],
                )
            )
        db.commit()
    finally:
        db.close()


# ==================== メインフロー ====================

def _indexing_sync(doc_id: str) -> int:
    """同期版本体。BackgroundTasks のスレッドで動かす。戻り値: 格納したチャンク数。"""
    db = SessionLocal()
    try:
        doc = db.query(Document).filter(Document.id == doc_id).first()
        if not doc:
            raise ValueError(f"Document not found: {doc_id}")

        file_path = doc.file_path
        source_type = doc.source_type
        title = doc.title or ""
        author = doc.author or ""
        year = doc.year or 0

        # OCR 文書の補正テキスト（再インデックス時はこちらを優先）
        corrected_text: Optional[str] = None
        if source_type == "ocr":
            existing_ocr = (
                db.query(OcrResult).filter(OcrResult.document_id == doc_id).first()
            )
            if existing_ocr and existing_ocr.corrected_text:
                corrected_text = existing_ocr.corrected_text
    finally:
        db.close()

    print(f"[pipeline] extracting file: {file_path}")

    if source_type == "ocr" and corrected_text:
        # 既存 OCR の補正テキストをそのまま markdown として使用する。
        # blocks は既に保存済みなので OCR は再実行しない。
        markdown_text = corrected_text
        ocr_result = None
        print(f"[pipeline] using corrected_text from ocr_results ({len(markdown_text)} chars)")
    else:
        markdown_text, ocr_result = _extract_to_markdown(source_type, file_path)
    print(f"[pipeline] markdown length: {len(markdown_text)} chars")

    if ocr_result is not None:
        print(
            f"[pipeline] ocr blocks: {len(ocr_result['blocks'])}, "
            f"avg_confidence: {ocr_result['avg_confidence']:.3f}, "
            f"low_conf_count: {ocr_result['low_conf_count']}"
        )
        _save_ocr_result(doc_id, ocr_result)

    print("[pipeline] chunking...")
    chunks = chunk_markdown(markdown_text)
    print(f"[pipeline] chunks before filter: {len(chunks)}")

    # Step 5: 品質フィルタリング（最低文字数 + 同一文書内重複除去）
    filtered = [c for c in chunks if len(c.page_content) >= settings.min_chunk_length]
    seen: set = set()
    deduped: list = []
    for chunk in filtered:
        key = (doc_id, chunk.page_content.strip())
        if key not in seen:
            seen.add(key)
            deduped.append(chunk)
    final_chunks = deduped
    print(f"[pipeline] chunks after filter: {len(final_chunks)}")

    if not final_chunks:
        raise ValueError("No chunks remaining after quality filtering")

    metadatas = [
        _build_metadata(c, doc_id, source_type, title, author, year)
        for c in final_chunks
    ]

    print(f"[pipeline] generating embeddings for {len(final_chunks)} chunks")
    embeddings = embed_chunks(final_chunks)

    # ChromaDB に格納する本文からはページマーカーを除去する
    for c in final_chunks:
        c.page_content = _strip_page_markers(c.page_content)

    print("[pipeline] storing to ChromaDB...")
    store_to_chroma(final_chunks, embeddings, metadatas)

    return len(final_chunks)


@with_retry(max_attempts=settings.pipeline_max_retry, base_delay=settings.pipeline_base_delay)
async def _run_indexing_pipeline_inner(doc_id: str) -> int:
    """リトライ対象の本体。同期処理を別スレッドにオフロードする。"""
    return await asyncio.to_thread(_indexing_sync, doc_id)


async def run_indexing_pipeline(doc_id: str) -> None:
    """
    文書インデックス化パイプラインのエントリーポイント。
    status の遷移: pending → indexing → indexed / error
    エラー発生時は error_message に詳細を記録して status=error にする。
    """
    print(f"[pipeline] start doc_id={doc_id}")

    db = SessionLocal()
    try:
        doc = db.query(Document).filter(Document.id == doc_id).first()
        if not doc:
            print(f"[pipeline] error: Document not found: {doc_id}")
            return
        doc.status = "indexing"
        doc.error_message = None
        db.commit()
    finally:
        db.close()

    try:
        chunk_count = await _run_indexing_pipeline_inner(doc_id)

        db = SessionLocal()
        try:
            doc = db.query(Document).filter(Document.id == doc_id).first()
            if doc:
                doc.status = "indexed"
                doc.chunk_count = chunk_count
                doc.error_message = None
                db.commit()
        finally:
            db.close()

        print(f"[pipeline] done. chunk_count={chunk_count}")
    except Exception as e:
        print(f"[pipeline] error: {e}")
        db = SessionLocal()
        try:
            doc = db.query(Document).filter(Document.id == doc_id).first()
            if doc:
                doc.status = "error"
                doc.error_message = str(e)[:1000]
                db.commit()
        finally:
            db.close()
